"""Extract text from PPTX and PDF files. Optional chapter detection for PPTX."""

import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pypdf import PdfReader

# Namespaces for PowerPoint section XML (p14 = 2010 main)
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"


def _normalize(text: str) -> str:
    """Normalize whitespace and strip."""
    if not text or not isinstance(text, str):
        return ""
    return re.sub(r"\s+", " ", text).strip()


def _slide_text(slide: Any) -> str:
    """Get all text from a slide as a single normalized string."""
    parts: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            parts.append(shape.text)
    return _normalize(" ".join(parts))


def _get_pptx_sections_native(prs: Presentation) -> list[tuple[str, list[int]]] | None:
    """
    Layer 1: Read native PowerPoint sections from presentation.xml.
    Returns list of (section_name, [slide_1based_indices]) or None if no sections.
    """
    root = prs.presentation.element
    # p:extLst > p:ext (with section uri) > p14:sectionLst > p14:section
    section_list = root.findall(f".//{{{NS_P14}}}section")
    if not section_list:
        return None
    # Build ordered list of slide IDs as in presentation (p:sldIdLst)
    sld_id_lst = root.find(f"{{{NS_P}}}sldIdLst")
    if sld_id_lst is None:
        return None
    order: list[str] = []
    for sld_id in sld_id_lst.findall(f"{{{NS_P}}}sldId"):
        id_val = sld_id.get("id")
        if id_val is not None:
            order.append(id_val)
    if not order:
        return None
    id_to_index = {sid: (i + 1) for i, sid in enumerate(order)}
    result: list[tuple[str, list[int]]] = []
    for sec in section_list:
        name = sec.get("name") or "Section"
        sld_id_lst_sec = sec.find(f"{{{NS_P14}}}sldIdLst")
        if sld_id_lst_sec is None:
            continue
        indices: list[int] = []
        for sld_id in sld_id_lst_sec.findall(f"{{{NS_P14}}}sldId"):
            id_val = sld_id.get("id")
            if id_val and id_val in id_to_index:
                indices.append(id_to_index[id_val])
        indices.sort()
        if indices:
            result.append((name, indices))
    return result if len(result) > 1 else None


# Heuristic: slide title matches chapter/section patterns
_CHAPTER_PATTERN = re.compile(
    r"^(chapter|week|unit|part|module|section|lecture)\s*(\d+|[ivxlcdm]+)?",
    re.IGNORECASE,
)


def _get_pptx_chapters_heuristic(prs: Presentation) -> list[tuple[str, list[tuple[int, str]]]] | None:
    """
    Layer 2: Detect chapter boundaries by section-header slides.
    A slide is a section start if it has little text and it matches Chapter N / Week N / etc.,
    or if it looks like a title-only slide (short text, could be section header).
    Returns list of (chapter_name, [(slide_index, text), ...]) or None.
    """
    items: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        text = _slide_text(slide)
        if text:
            items.append((i, text))
    if not items:
        return None
    # Find section-head slides: first line or full text matches chapter pattern, or very short
    section_starts: list[tuple[int, str]] = []  # (slide_idx, title for chapter name)
    for idx, text in items:
        first_line = text.split("\n")[0].strip() if text else ""
        if _CHAPTER_PATTERN.match(first_line) or _CHAPTER_PATTERN.match(text):
            section_starts.append((idx, first_line or text[:80]))
        elif len(text) < 120 and not any(c in text for c in [".", ":", "?"]):
            # Short title-like slide, likely a section header
            section_starts.append((idx, first_line or text[:80]))
    if not section_starts:
        return None
    # Build chapters: from start to next start (or end)
    chapters: list[tuple[str, list[tuple[int, str]]]] = []
    for k, (start_idx, title) in enumerate(section_starts):
        end_idx = section_starts[k + 1][0] if k + 1 < len(section_starts) else items[-1][0] + 1
        # Exclude section-header slide from content (boundary only)
        chunk_items = [(i, t) for i, t in items if start_idx < i < end_idx]
        if not chunk_items:
            continue
        chapter_name = _normalize(title) or f"Part {k + 1}"
        chapters.append((chapter_name, chunk_items))
    # Prepend slides before first section as "Introduction" or first chapter
    first_start = section_starts[0][0]
    intro = [(i, t) for i, t in items if i < first_start]
    if intro:
        chapters.insert(0, ("Introduction", intro))
    return chapters if len(chapters) >= 2 else None


def extract_pptx(path: str | Path) -> list[tuple[int, str]]:
    """Extract text from each slide of a PPTX file. Returns list of (slide_index, text)."""
    path = Path(path)
    if not path.suffix.lower() == ".pptx":
        raise ValueError(f"Expected .pptx file, got {path.suffix}")
    prs = Presentation(str(path))
    result: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        text = _slide_text(slide)
        if text:
            result.append((i, text))
    return result


def extract_pptx_with_chapters(path: str | Path) -> list[tuple[str, list[tuple[int, str]]]]:
    """
    Extract PPTX with chapter detection. Returns list of (chapter_name, [(slide_index, text), ...]).
    Tries native sections first, then heuristic. If no chapters found, returns one chapter with all slides.
    """
    path = Path(path)
    if path.suffix.lower() != ".pptx":
        raise ValueError(f"Expected .pptx file, got {path.suffix}")
    prs = Presentation(str(path))
    items = [(i, _slide_text(slide)) for i, slide in enumerate(prs.slides, start=1)]
    items = [(i, t) for i, t in items if t]

    sections = _get_pptx_sections_native(prs)
    if sections:
        # Map (name, [indices]) to (name, [(idx, text)])
        result: list[tuple[str, list[tuple[int, str]]]] = []
        for name, indices in sections:
            chapter_items = [(i, next((t for i2, t in items if i2 == i), "")) for i in indices]
            chapter_items = [(i, t) for i, t in chapter_items if t]
            if chapter_items:
                result.append((name, chapter_items))
        if result:
            return result

    heuristic = _get_pptx_chapters_heuristic(prs)
    if heuristic:
        return heuristic

    return [("Full deck", items)] if items else []


def extract_pdf(path: str | Path) -> list[tuple[int, str]]:
    """Extract text from each page of a PDF. Returns list of (page_number, text)."""
    path = Path(path)
    if not path.suffix.lower() == ".pdf":
        raise ValueError(f"Expected .pdf file, got {path.suffix}")
    reader = PdfReader(str(path))
    result: list[tuple[int, str]] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        text = _normalize(text or "")
        if text:
            result.append((i, text))
    return result


def extract(path: str | Path) -> list[tuple[int, str]]:
    """
    Extract text from a PPTX or PDF file. Infers format from extension.
    Returns list of (slide_or_page_index, text). Empty slides/pages are omitted.
    """
    path = Path(path)
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".pdf":
        return extract_pdf(path)
    raise ValueError(f"Unsupported format: {suffix}. Use .pptx or .pdf")
